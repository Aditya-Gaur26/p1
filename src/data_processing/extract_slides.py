"""
Extract text from PowerPoint slides with enhanced preprocessing
"""

import os
from pathlib import Path
from pptx import Presentation
from typing import List, Dict
import sys

sys.path.append(str(Path(__file__).parent.parent.parent))
from src.utils.config import config
from src.utils.helpers import save_json, clean_text
from src.utils.text_preprocessing import AdvancedTextCleaner, clean_slide_text


def extract_text_from_slide(slide) -> str:
    """Extract text from a single slide with enhanced cleaning"""
    text_parts = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                # Apply slide-specific cleaning
                text = clean_slide_text(text)
                if text:  # Only add non-empty after cleaning
                    text_parts.append(text)
    
    return "\n".join(text_parts)


def extract_from_pptx(filepath: Path) -> Dict:
    """Extract content from PowerPoint file with structure detection"""
    try:
        prs = Presentation(filepath)
        cleaner = AdvancedTextCleaner()
        
        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = extract_text_from_slide(slide)
            if slide_text:
                # Extract structured content from slide
                structured = cleaner.extract_structured_content(slide_text)
                
                slides_content.append({
                    "slide_number": i,
                    "content": slide_text,
                    "has_bullets": len(structured.get('bullet_points', [])) > 0,
                    "has_code": cleaner.is_code_block(slide_text),
                    "headers": structured.get('headers', [])
                })
        
        return {
            "filename": filepath.name,
            "total_slides": len(prs.slides),
            "slides": slides_content,
            "text": "\n\n".join([s["content"] for s in slides_content])
        }
    
    except Exception as e:
        print(f"❌ Error processing {filepath.name}: {str(e)}")
        return None


def process_all_slides(input_dir: Path, output_dir: Path):
    """Process all PowerPoint files in directory"""
    input_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Find all PPTX files
    pptx_files = list(input_dir.glob("*.pptx")) + list(input_dir.glob("*.ppt"))
    
    if not pptx_files:
        print(f"⚠️  No PowerPoint files found in {input_dir}")
        print(f"Please add .pptx or .ppt files to {input_dir}")
        return
    
    print(f"📊 Found {len(pptx_files)} PowerPoint files")
    
    all_extracted = []
    
    for pptx_file in pptx_files:
        print(f"Processing: {pptx_file.name}...")
        extracted = extract_from_pptx(pptx_file)
        
        if extracted:
            all_extracted.append(extracted)
            
            # Save individual file
            output_file = output_dir / f"{pptx_file.stem}_extracted.json"
            save_json(extracted, output_file)
            print(f"  ✓ Extracted {extracted['total_slides']} slides")
    
    # Save combined output
    if all_extracted:
        combined_output = output_dir / "all_slides_combined.json"
        save_json({"files": all_extracted, "total_files": len(all_extracted)}, combined_output)
        print(f"\n✅ Successfully processed {len(all_extracted)} files")
        print(f"📁 Output saved to: {output_dir}")
    else:
        print("❌ No files were successfully processed")


def main():
    """Main execution function"""
    input_dir = config.data_dir / "raw" / "slides"
    output_dir = config.data_dir / "processed" / "slides"
    
    print("=" * 60)
    print("PowerPoint Slide Extraction".center(60))
    print("=" * 60)
    
    process_all_slides(input_dir, output_dir)


if __name__ == "__main__":
    main()
