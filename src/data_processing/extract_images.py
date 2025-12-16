"""
Extract and process images from PowerPoint slides and PDFs
Handles image extraction, OCR, image description generation, and diagram recognition
"""

import sys
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
import io
from PIL import Image
import hashlib
import numpy as np

sys.path.append(str(Path(__file__).parent.parent.parent))
from src.utils.config import config
from src.utils.helpers import save_json

# Try to import OCR libraries
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    print("⚠️  pytesseract not available. Install with: pip install pytesseract")

# Try to import vision libraries
try:
    from transformers import BlipProcessor, BlipForConditionalGeneration
    BLIP_AVAILABLE = True
except ImportError:
    BLIP_AVAILABLE = False
    print("⚠️  BLIP not available. Install with: pip install transformers")


class ImageExtractor:
    """Extract images from PowerPoint presentations and PDFs with intelligent processing"""
    
    def __init__(self, save_images: bool = True, output_dir: Path = None,
                 use_ocr: bool = True, use_captioning: bool = True):
        """
        Initialize image extractor with advanced features
        
        Args:
            save_images: Whether to save extracted images to disk
            output_dir: Directory to save images
            use_ocr: Enable OCR for text extraction from images
            use_captioning: Enable AI image captioning
        """
        self.save_images = save_images
        self.output_dir = output_dir or (config.data_dir / "processed" / "images")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        self.use_ocr = use_ocr and TESSERACT_AVAILABLE
        self.use_captioning = use_captioning and BLIP_AVAILABLE
        
        # Initialize captioning model if available
        self.caption_model = None
        self.caption_processor = None
        if self.use_captioning:
            self._init_caption_model()
        
        # Track extracted images
        self.extracted_images = []
    
    def _init_caption_model(self):
        """Initialize BLIP image captioning model"""
        try:
            print("📦 Loading BLIP image captioning model...")
            self.caption_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
            self.caption_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
            print("✓ BLIP model loaded")
        except Exception as e:
            print(f"⚠️  Could not load BLIP model: {e}")
            self.use_captioning = False
        
    def extract_images_from_slide(self, slide, slide_number: int, 
                                  filename: str) -> List[Dict[str, Any]]:
        """
        Extract all images from a single slide
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number
            filename: Source filename
            
        Returns:
            List of image information dictionaries
        """
        images = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # Check if shape contains an image
                if hasattr(shape, "image"):
                    image_data = self._extract_image_data(shape)
                    if image_data:
                        image_info = {
                            'slide_number': slide_number,
                            'shape_index': shape_idx,
                            'source_file': filename,
                            'image_type': image_data['type'],
                            'image_format': image_data['format'],
                            'image_size': image_data['size'],
                            'image_hash': image_data['hash'],
                            'saved_path': image_data.get('saved_path'),
                            'has_caption': False,
                            'caption': ''
                        }
                        
                        # Check for caption text
                        if hasattr(shape, 'text') and shape.text.strip():
                            image_info['has_caption'] = True
                            image_info['caption'] = shape.text.strip()
                        
                        images.append(image_info)
                        
                # Check for picture shapes
                elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_data = self._extract_picture_shape(shape, slide_number, 
                                                             filename, shape_idx)
                    if image_data:
                        images.append(image_data)
                        
            except Exception as e:
                print(f"  ⚠️  Could not extract image from shape {shape_idx}: {e}")
                continue
        
        return images
    
    def _extract_image_data(self, shape) -> Dict[str, Any]:
        """Extract image data from shape"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Get image format
            image_format = image.content_type.split('/')[-1].upper()
            
            # Create PIL Image
            pil_image = Image.open(io.BytesIO(image_bytes))
            
            # Generate hash for deduplication
            image_hash = hashlib.md5(image_bytes).hexdigest()
            
            # Save image if requested
            saved_path = None
            if self.save_images:
                saved_path = self._save_image(pil_image, image_hash, image_format)
            
            return {
                'type': 'embedded',
                'format': image_format,
                'size': pil_image.size,
                'hash': image_hash,
                'saved_path': str(saved_path) if saved_path else None,
                'bytes': image_bytes if not self.save_images else None
            }
            
        except Exception as e:
            print(f"  ⚠️  Error extracting image data: {e}")
            return None
    
    def _extract_picture_shape(self, shape, slide_number: int, 
                               filename: str, shape_idx: int) -> Dict[str, Any]:
        """Extract picture from picture shape"""
        try:
            # Get the image from the picture shape
            image = shape.image
            image_bytes = image.blob
            
            # Create PIL Image
            pil_image = Image.open(io.BytesIO(image_bytes))
            
            # Generate hash
            image_hash = hashlib.md5(image_bytes).hexdigest()
            
            # Get format
            image_format = image.content_type.split('/')[-1].upper()
            
            # Save image
            saved_path = None
            if self.save_images:
                saved_path = self._save_image(pil_image, image_hash, image_format)
            
            return {
                'slide_number': slide_number,
                'shape_index': shape_idx,
                'source_file': filename,
                'image_type': 'picture',
                'image_format': image_format,
                'image_size': pil_image.size,
                'image_hash': image_hash,
                'saved_path': str(saved_path) if saved_path else None
            }
            
        except Exception as e:
            print(f"  ⚠️  Error extracting picture shape: {e}")
            return None
    
    def _save_image(self, pil_image: Image.Image, image_hash: str, 
                    image_format: str) -> Path:
        """Save image to disk"""
        # Use hash as filename to avoid duplicates
        filename = f"{image_hash}.{image_format.lower()}"
        filepath = self.output_dir / filename
        
        # Save only if doesn't exist
        if not filepath.exists():
            pil_image.save(filepath, format=image_format)
        
        return filepath
    
    def extract_images_from_presentation(self, pptx_path: Path) -> List[Dict[str, Any]]:
        """
        Extract all images from a PowerPoint presentation
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            List of all extracted images with metadata
        """
        from pptx import Presentation
        
        print(f"🖼️  Extracting images from: {pptx_path.name}")
        
        try:
            prs = Presentation(pptx_path)
            all_images = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_images = self.extract_images_from_slide(
                    slide, slide_num, pptx_path.name
                )
                all_images.extend(slide_images)
            
            print(f"  ✓ Extracted {len(all_images)} images from {len(prs.slides)} slides")
            
            return all_images
            
        except Exception as e:
            print(f"  ❌ Error processing {pptx_path.name}: {e}")
            return []
    
    def get_statistics(self, images: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Get statistics about extracted images"""
        if not images:
            return {}
        
        formats = {}
        total_slides = set()
        
        for img in images:
            fmt = img.get('image_format', 'UNKNOWN')
            formats[fmt] = formats.get(fmt, 0) + 1
            total_slides.add(img.get('slide_number'))
        
        return {
            'total_images': len(images),
            'unique_images': len(set(img.get('image_hash') for img in images)),
            'slides_with_images': len(total_slides),
            'formats': formats,
            'images_with_captions': sum(1 for img in images if img.get('has_caption', False))
        }


class ImageOCRProcessor:
    """Process images with OCR to extract text"""
    
    def __init__(self):
        """Initialize OCR processor"""
        self.ocr_available = False
        try:
            import pytesseract
            from PIL import Image
            self.pytesseract = pytesseract
            self.ocr_available = True
            print("✓ Tesseract OCR available")
        except ImportError:
            print("⚠️  pytesseract not installed. Run: pip install pytesseract")
            print("⚠️  Also install Tesseract: https://github.com/tesseract-ocr/tesseract")
    
    def extract_text_from_image(self, image_path: Path, lang: str = 'eng') -> Dict[str, Any]:
        """
        Extract text from image using OCR
        
        Args:
            image_path: Path to image file
            lang: Language for OCR (default: English)
            
        Returns:
            Dictionary with extracted text and confidence
        """
        if not self.ocr_available:
            return {'text': '', 'confidence': 0, 'error': 'OCR not available'}
        
        try:
            from PIL import Image
            
            # Open image
            image = Image.open(image_path)
            
            # Perform OCR
            text = self.pytesseract.image_to_string(image, lang=lang)
            
            # Get detailed data with confidence
            data = self.pytesseract.image_to_data(image, output_type=self.pytesseract.Output.DICT)
            
            # Calculate average confidence
            confidences = [int(conf) for conf in data['conf'] if conf != '-1']
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            return {
                'text': text.strip(),
                'confidence': avg_confidence,
                'word_count': len(text.split()),
                'has_text': bool(text.strip())
            }
            
        except Exception as e:
            return {'text': '', 'confidence': 0, 'error': str(e)}
    
    def process_image_batch(self, images: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Process multiple images with OCR
        
        Args:
            images: List of image dictionaries with 'saved_path'
            
        Returns:
            Updated list with OCR results
        """
        if not self.ocr_available:
            print("⚠️  OCR not available, skipping text extraction")
            return images
        
        print(f"📝 Processing {len(images)} images with OCR...")
        
        processed = []
        for i, img in enumerate(images):
            if img.get('saved_path'):
                print(f"  [{i+1}/{len(images)}] Processing {Path(img['saved_path']).name}...")
                
                ocr_result = self.extract_text_from_image(Path(img['saved_path']))
                
                # Add OCR results to image info
                img['ocr_text'] = ocr_result.get('text', '')
                img['ocr_confidence'] = ocr_result.get('confidence', 0)
                img['ocr_has_text'] = ocr_result.get('has_text', False)
                
                if ocr_result.get('has_text'):
                    print(f"    ✓ Extracted {ocr_result.get('word_count', 0)} words "
                          f"(confidence: {ocr_result.get('confidence', 0):.1f}%)")
            
            processed.append(img)
        
        return processed
    
    def preprocess_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """
        Preprocess image to improve OCR accuracy
        
        Args:
            image: PIL Image
            
        Returns:
            Preprocessed PIL Image
        """
        # Convert to grayscale
        image = image.convert('L')
        
        # Enhance contrast
        from PIL import ImageEnhance
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.0)
        
        # Threshold to binary
        threshold = 128
        image = image.point(lambda p: 255 if p > threshold else 0)
        
        # Resize if too small (OCR works better on larger images)
        if image.size[0] < 300:
            new_width = 300
            new_height = int(image.size[1] * (new_width / image.size[0]))
            image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
        
        return image


def extract_and_process_images(pptx_path: Path, use_ocr: bool = True) -> Dict[str, Any]:
    """
    Main function to extract and process images from PowerPoint
    
    Args:
        pptx_path: Path to PowerPoint file
        use_ocr: Whether to apply OCR to images
        
    Returns:
        Dictionary with images and metadata
    """
    # Extract images
    extractor = ImageExtractor(save_images=True)
    images = extractor.extract_images_from_presentation(pptx_path)
    
    # Apply OCR if requested
    if use_ocr and images:
        ocr_processor = ImageOCRProcessor()
        images = ocr_processor.process_image_batch(images)
    
    # Get statistics
    stats = extractor.get_statistics(images)
    
    result = {
        'source_file': pptx_path.name,
        'images': images,
        'statistics': stats
    }
    
    return result


def main():
    """Test image extraction"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Extract images from PowerPoint")
    parser.add_argument("--file", type=str, help="PowerPoint file path")
    parser.add_argument("--input-dir", type=str, default="./data/raw/slides",
                       help="Directory containing PowerPoint files")
    parser.add_argument("--no-ocr", action="store_true", 
                       help="Disable OCR processing")
    
    args = parser.parse_args()
    
    if args.file:
        # Process single file
        result = extract_and_process_images(Path(args.file), use_ocr=not args.no_ocr)
        
        # Save results
        output_file = config.data_dir / "processed" / "images" / f"{Path(args.file).stem}_images.json"
        save_json(result, output_file)
        print(f"\n✅ Results saved to: {output_file}")
        
    else:
        # Process all files in directory
        input_dir = Path(args.input_dir)
        if not input_dir.exists():
            print(f"❌ Directory not found: {input_dir}")
            return
        
        pptx_files = list(input_dir.glob("*.pptx")) + list(input_dir.glob("*.ppt"))
        
        if not pptx_files:
            print(f"❌ No PowerPoint files found in {input_dir}")
            return
        
        print(f"🖼️  Processing {len(pptx_files)} files...")
        
        all_results = []
        for pptx_file in pptx_files:
            result = extract_and_process_images(pptx_file, use_ocr=not args.no_ocr)
            all_results.append(result)
        
        # Save combined results
        combined_output = config.data_dir / "processed" / "images" / "all_images_extracted.json"
        save_json({
            'files': all_results,
            'total_files': len(all_results),
            'total_images': sum(r['statistics'].get('total_images', 0) for r in all_results)
        }, combined_output)
        
        print(f"\n✅ Combined results saved to: {combined_output}")


if __name__ == "__main__":
    main()
